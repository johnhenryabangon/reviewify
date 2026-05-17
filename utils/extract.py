"""Smarter extraction from PDF / PPT / PPTX.

Returns a list of sections per file:
    [{
        "heading": str,
        "subheading": str,
        "bullets": [str, ...],
        "body": str,
        "keywords": [str, ...],
        "formulas": [str, ...],
        "images": [str, ...],   # web paths under /static/uploads/...
    }, ...]

Image extraction:
  - PPTX: every PICTURE shape on a slide is associated with that slide's section.
  - PDF:  PyMuPDF extracts embedded raster images per page; tiny / decorative
          images are filtered out by minimum dimensions.
"""
from __future__ import annotations

import io
import re
import hashlib
from pathlib import Path
from typing import List, Dict, Tuple, Optional


# -------------------------------------------------------------------- public
def extract_document(path: Path, image_dir: Optional[Path] = None,
                     image_url_prefix: str = "") -> List[Dict]:
    """Extract structured sections + images.

    image_dir:        filesystem dir to write extracted images into.
    image_url_prefix: URL prefix that maps to image_dir for the browser
                      (e.g. "/static/uploads/<batch>/images").
    """
    ext = path.suffix.lower()
    if ext == ".pdf":
        sections = _extract_pdf(path, image_dir, image_url_prefix)
    elif ext in {".pptx", ".ppt"}:
        sections = _extract_pptx(path, image_dir, image_url_prefix)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    sections = _dedupe_repeated_lines(sections)
    sections = _enrich(sections)
    return sections


# --------------------------------------------------------------------- PPTX
def _extract_pptx(path: Path, image_dir: Optional[Path],
                  url_prefix: str) -> List[Dict]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    sections: List[Dict] = []
    for i, slide in enumerate(prs.slides, start=1):
        heading = ""
        subheading = ""
        bullets: List[str] = []
        body_parts: List[str] = []
        bold_terms: List[str] = []
        images: List[str] = []

        if slide.shapes.title and slide.shapes.title.has_text_frame:
            heading = (slide.shapes.title.text or "").strip()

        for shape in slide.shapes:
            # ---- pictures ----
            if image_dir is not None and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    ext_ = (shape.image.ext or "png").lower()
                    url = _save_image(blob, ext_, image_dir, url_prefix,
                                      min_w=140, min_h=140)
                    if url:
                        images.append(url)
                except Exception:
                    pass
                continue

            if shape == slide.shapes.title or not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if not text:
                    continue
                # Track bold phrases as keyword candidates
                for run in para.runs:
                    if run.font.bold and run.text and len(run.text.strip()) > 2:
                        bold_terms.append(run.text.strip(" :;.,-"))
                if not heading:
                    heading = text
                    continue
                if not subheading and len(text) < 90 and not bullets and not body_parts:
                    subheading = text
                    continue
                if len(text) < 220:
                    bullets.append(text)
                else:
                    body_parts.append(text)

        if not heading:
            heading = f"Slide {i}"
        if bullets or body_parts or images:
            sections.append({
                "heading": heading,
                "subheading": subheading,
                "bullets": bullets,
                "body": "\n".join(body_parts).strip(),
                "keywords": _uniq(bold_terms)[:8],
                "formulas": [],
                "images": images,
            })
    return sections


# ---------------------------------------------------------------------- PDF
def _extract_pdf(path: Path, image_dir: Optional[Path],
                 url_prefix: str) -> List[Dict]:
    import pdfplumber

    # ---- text first, with per-page tracking ----
    page_sections: Dict[int, List[Dict]] = {}
    current = None
    current_page = 1

    with pdfplumber.open(str(path)) as pdf:
        for pno, page in enumerate(pdf.pages, start=1):
            current_page = pno
            text = page.extract_text() or ""
            for raw in text.splitlines():
                line = raw.strip()
                if not line:
                    continue
                if _looks_like_heading(line):
                    if current and (current["bullets"] or current["body"]):
                        page_sections.setdefault(current["_page"], []).append(current)
                    current = _new_sec(line, pno)
                else:
                    if current is None:
                        current = _new_sec("Introduction", pno)
                    if _looks_like_bullet(line):
                        current["bullets"].append(_strip_bullet(line))
                    else:
                        current["body"] = (current["body"] + " " + line).strip()
    if current and (current["bullets"] or current["body"]):
        page_sections.setdefault(current["_page"], []).append(current)

    # ---- images via PyMuPDF, attached to page's first section ----
    if image_dir is not None:
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(str(path))
            for pno, page in enumerate(doc, start=1):
                imgs_on_page: List[str] = []
                for img in page.get_images(full=True):
                    xref = img[0]
                    try:
                        pix = fitz.Pixmap(doc, xref)
                        if pix.n - pix.alpha >= 4:  # CMYK -> RGB
                            pix = fitz.Pixmap(fitz.csRGB, pix)
                        if pix.width < 140 or pix.height < 140:
                            pix = None
                            continue
                        blob = pix.tobytes("png")
                        url = _save_image(blob, "png", image_dir, url_prefix,
                                          min_w=140, min_h=140)
                        if url:
                            imgs_on_page.append(url)
                        pix = None
                    except Exception:
                        continue
                if imgs_on_page:
                    bucket = page_sections.setdefault(pno, [
                        _new_sec(f"Page {pno}", pno)
                    ])
                    bucket[0]["images"].extend(imgs_on_page)
            doc.close()
        except Exception as e:
            print(f"[extract] PDF image extraction unavailable: {e}")

    # ---- flatten in page order ----
    sections: List[Dict] = []
    for pno in sorted(page_sections.keys()):
        for s in page_sections[pno]:
            s.pop("_page", None)
            sections.append(s)

    if not sections:
        sections = [_new_sec("Document", 1)]
        for s in sections:
            s.pop("_page", None)
    return sections


# ----------------------------------------------------------------- helpers
def _new_sec(heading: str, page: int) -> Dict:
    return {
        "_page": page,
        "heading": heading,
        "subheading": "",
        "bullets": [],
        "body": "",
        "keywords": [],
        "formulas": [],
        "images": [],
    }


def _looks_like_heading(line: str) -> bool:
    if len(line) > 90 or len(line) < 3:
        return False
    letters = [c for c in line if c.isalpha()]
    if not letters:
        return False
    upper_ratio = sum(1 for c in letters if c.isupper()) / len(letters)
    if upper_ratio > 0.7 and len(line.split()) <= 12:
        return True
    lower = line.lower()
    if lower.startswith(("chapter ", "lesson ", "module ", "unit ",
                         "topic ", "section ", "part ")):
        return True
    if line[:2].rstrip(".").isdigit() and len(line) < 80 and line.endswith(
            tuple("abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ")):
        return True
    return False


def _looks_like_bullet(line: str) -> bool:
    if line[:2] in {"- ", "• ", "* ", "▪ ", "› "}:
        return True
    return bool(re.match(r"^\d{1,2}[.)]\s+\S", line)) and len(line) < 220


def _strip_bullet(line: str) -> str:
    for prefix in ("- ", "• ", "* ", "▪ ", "› "):
        if line.startswith(prefix):
            return line[len(prefix):].strip()
    m = re.match(r"^\d{1,2}[.)]\s+(.*)$", line)
    return m.group(1).strip() if m else line


def _dedupe_repeated_lines(sections: List[Dict]) -> List[Dict]:
    """Drop bullets that repeat verbatim across many sections (slide footers,
    course code, instructor name, etc.)."""
    counts: Dict[str, int] = {}
    for s in sections:
        for b in s["bullets"]:
            k = b.strip().lower()
            if 4 < len(k) < 80:
                counts[k] = counts.get(k, 0) + 1
    threshold = max(3, len(sections) // 3)
    noisy = {k for k, c in counts.items() if c >= threshold}
    for s in sections:
        s["bullets"] = [b for b in s["bullets"]
                        if b.strip().lower() not in noisy]
        # collapse duplicate bullets within a section
        seen = set()
        uniq = []
        for b in s["bullets"]:
            kk = b.strip().lower()
            if kk in seen:
                continue
            seen.add(kk)
            uniq.append(b)
        s["bullets"] = uniq
    return sections


_FORMULA_RE = re.compile(
    r"[A-Za-z0-9_\)\]]\s*[=≈≠≤≥<>]\s*[^,;]{1,80}"   # x = something
    r"|[∑∏∫√±×÷→⇒⇔αβγδθλμπσφω]"                     # math symbol
    r"|\bO\(\s*[nlogN0-9 ^*+\-]+\s*\)"               # O(n log n)
)


def _enrich(sections: List[Dict]) -> List[Dict]:
    for s in sections:
        text_pool = " ".join([s["heading"], s["subheading"], s["body"]] + s["bullets"])
        # formulas
        formulas = []
        for line in [s["body"]] + s["bullets"]:
            for m in _FORMULA_RE.findall(line or ""):
                f = m.strip()
                if 3 <= len(f) <= 120 and f not in formulas:
                    formulas.append(f)
        s["formulas"] = formulas[:6]
        # keyword candidates: capitalized multi-word terms + bold terms
        caps = re.findall(r"\b([A-Z][a-zA-Z]{2,}(?:\s+[A-Z][a-zA-Z]+){0,3})\b",
                          text_pool)
        kws = _uniq(s["keywords"] + caps)
        # filter trivial
        kws = [k for k in kws if len(k) > 3 and k.lower() not in _STOP_KW]
        s["keywords"] = kws[:8]
    return sections


_STOP_KW = {"the", "this", "that", "these", "those", "introduction",
            "overview", "summary", "agenda", "outline", "objectives",
            "references", "questions", "thank", "thanks"}


def _uniq(xs: List[str]) -> List[str]:
    seen = set()
    out = []
    for x in xs:
        k = x.lower().strip()
        if not k or k in seen:
            continue
        seen.add(k)
        out.append(x.strip())
    return out


# -------------------------------------------------------------- image utils
def _save_image(blob: bytes, ext: str, image_dir: Path, url_prefix: str,
                min_w: int = 120, min_h: int = 120) -> Optional[str]:
    """Write image blob to disk, dedupe by hash, filter tiny/uniform images.
    Returns a web URL or None if the image was rejected."""
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(blob))
        w, h = img.size
        if w < min_w or h < min_h:
            return None
        # reject near-uniform (likely decorative bar / background)
        try:
            small = img.convert("RGB").resize((16, 16))
            colors = small.getcolors(maxcolors=256) or []
            if colors:
                colors.sort(reverse=True)
                if colors[0][0] > 220:  # >85% of pixels are one color
                    return None
        except Exception:
            pass

        h_hash = hashlib.md5(blob).hexdigest()[:12]
        ext = (ext or "png").lstrip(".")
        if ext not in {"png", "jpg", "jpeg", "gif", "webp", "bmp"}:
            ext = "png"
        fname = f"{h_hash}.{ext}"
        out_path = image_dir / fname
        if not out_path.exists():
            image_dir.mkdir(parents=True, exist_ok=True)
            with open(out_path, "wb") as fh:
                fh.write(blob)
        return f"{url_prefix.rstrip('/')}/{fname}"
    except Exception:
        return None
